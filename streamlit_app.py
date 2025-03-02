import io
import streamlit as st
from pptx import Presentation
from googletrans import Translator

st.title("Googletrans PowerPoint Translator")

def translate_text(text, dest_language):
    """Translate text using googletrans (unofficial)."""
    translator = Translator()
    result = translator.translate(text, dest=dest_language)
    return result.text

def process_pptx_bytes(in_bytes):
    """
    1) Reads the PPTX file in memory.
    2) For each slide, grabs the notes text.
    3) Translates the notes to Amharic & Swedish.
    4) Overwrites the existing notes with the translated versions (or you can add shapes, etc.).
    5) Returns a BytesIO of the new PPTX.
    """
    prs = Presentation(io.BytesIO(in_bytes))

    for slide in prs.slides:
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                amharic = translate_text(notes_text, 'am')
                swedish = translate_text(notes_text, 'sv')

                # Overwrite the notes (simple approach)
                slide.notes_slide.notes_text_frame.text = (
                    f"Amharic:\n{amharic}\n\nSwedish:\n{swedish}"
                )

    out_bytes = io.BytesIO()
    prs.save(out_bytes)
    out_bytes.seek(0)
    return out_bytes

uploaded_file = st.file_uploader("Upload a .pptx file", type=["pptx"])
if uploaded_file is not None:
    st.write("Translating... please wait.")
    # Read the uploaded file into memory
    ppt_in_memory = uploaded_file.read()

    # Process (translate) the PPTX in memory
    translated_ppt = process_pptx_bytes(ppt_in_memory)

    st.success("Translation complete! Click below to download your PPTX.")
    st.download_button(
        label="Download Translated PPTX",
        data=translated_ppt,
        file_name="translated_output.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
